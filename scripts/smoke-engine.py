"""Validate conversions through the frozen worker without importing feature sources."""
import argparse
import json
from pathlib import Path
import subprocess
import tempfile

from PIL import Image
from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas


def require(condition, message):
    if not condition:
        raise RuntimeError(message)


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--root", type=Path, default=Path(__file__).resolve().parent.parent,
                        help="Checkout containing the worker binary and feature registry")
    args = parser.parse_args()
    root = args.root.resolve()
    executable = root / "src-tauri/binaries/pdf-worker-x86_64-pc-windows-msvc.exe"
    require(executable.is_file(), f"Packaged worker missing: {executable}")
    available = {path.stem for path in (root / "engine/features").glob("*.py") if not path.name.startswith("_")}
    completed = []
    with tempfile.TemporaryDirectory(prefix="pdf-worker-smoke-") as temp:
        # Windows CI can expose TEMP with an 8.3 alias (RUNNER~1). Match the
        # canonical paths returned by the worker before checking confinement.
        folder = Path(temp).resolve()
        source = folder / "entrada.pdf"
        picture = folder / "imagen.png"
        form = folder / "formulario.pdf"
        encrypted = folder / "protegido.pdf"
        Image.new("RGB", (60, 40), "red").save(picture)
        document = canvas.Canvas(str(source), pagesize=(300, 400))
        document.drawString(20, 350, "Frozen worker page one")
        document.drawImage(str(picture), 20, 200, 120, 80)
        document.showPage()
        document.drawString(20, 350, "Frozen worker page two")
        document.save()
        document = canvas.Canvas(str(form), pagesize=(300, 400))
        document.acroForm.textfield(name="nombre", x=20, y=200, width=200, height=30)
        document.showPage()
        document.save()
        writer = PdfWriter(clone_from=source)
        writer.encrypt("smoke-password", algorithm="AES-256")
        writer.write(encrypted)
        originals = {path: path.read_bytes() for path in (source, picture, form, encrypted)}

        def invoke(feature, inputs, options=None):
            request = {"feature": feature, "inputs": [str(path) for path in inputs],
                       "output_dir": str(folder), "options": options or {}}
            # Running away from the checkout prevents accidental source imports.
            result = subprocess.run([str(executable)], input=json.dumps(request).encode("utf-8"),
                                    capture_output=True, timeout=90, cwd=folder,
                                    creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
            require(result.returncode == 0, f"{feature}: worker crashed: {result.stderr.decode(errors='replace')}")
            try:
                response = json.loads(result.stdout)
            except (ValueError, UnicodeDecodeError) as error:
                raise RuntimeError(f"{feature}: invalid worker JSON: {result.stdout[:500]!r}") from error
            return response

        invalid = invoke("unknown", [source])
        require(invalid.get("ok") is False and invalid.get("outputs") == [], "Frozen registry accepted unknown feature")
        completed.append("registry rejection")

        def run(feature, inputs, options=None):
            response = invoke(feature, inputs, options)
            require(response.get("ok") is True, f"{feature}: {response}")
            paths = [Path(path).resolve() for path in response.get("outputs", [])]
            require(bool(paths), f"{feature}: no output files")
            for path in paths:
                require(path.is_file() and path.stat().st_size > 0, f"{feature}: empty or missing output {path}")
                require(path.is_relative_to(folder) and path not in originals, f"{feature}: invalid output location {path}")
            completed.append(feature)
            print(f"Frozen worker: {feature} executed", flush=True)
            return paths

        def pdf_text(path):
            return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)

        if "merge" in available:
            result = run("merge", [source, source])[0]
            require(len(PdfReader(result).pages) == 4, "merge: wrong page count")
            require(pdf_text(result).count("Frozen worker page one") == 2, "merge: text lost")

        if "pdf_to_jpg" in available:
            results = run("pdf_to_jpg", [source], {"dpi": 72})
            require(len(results) == 2, "pdf_to_jpg: wrong number of rendered pages")
            with Image.open(results[0]) as rendered:
                require(rendered.format == "JPEG" and rendered.size == (300, 400), "pdf_to_jpg: incorrect rendering")
                require(rendered.getpixel((40, 140))[0] > 230 and rendered.getpixel((40, 140))[1] < 30,
                        "pdf_to_jpg: embedded picture missing in rendered page")
            results = run("pdf_to_jpg", [source], {"mode": "Extraer imágenes"})
            require(len(results) == 1, "pdf_to_jpg: embedded image extraction failed")
            with Image.open(results[0]) as extracted:
                require(extracted.size == (60, 40), "pdf_to_jpg: extracted image resolution changed")

        if "pdf_to_word" in available:
            from docx import Document
            result = run("pdf_to_word", [source])[0]
            text = "\n".join(paragraph.text for paragraph in Document(result).paragraphs)
            require("Frozen worker page one" in text and "Frozen worker page two" in text, "pdf_to_word: missing editable text")

        if "pdf_to_powerpoint" in available:
            from pptx import Presentation
            presentation = Presentation(run("pdf_to_powerpoint", [source])[0])
            require(len(presentation.slides) == 2, "pdf_to_powerpoint: missing slide")
            require(all(len(slide.shapes) == 1 and slide.shapes[0].image.size == (600, 800)
                        for slide in presentation.slides), "pdf_to_powerpoint: page images missing")

        if "pdf_to_excel" in available:
            from openpyxl import load_workbook
            workbook = load_workbook(run("pdf_to_excel", [source])[0])
            try:
                require(len(workbook.worksheets) == 2, "pdf_to_excel: missing sheets")
                require(workbook.worksheets[0]["A1"].value == "Frozen worker page one", "pdf_to_excel: missing extracted text")
            finally:
                workbook.close()

        if "watermark" in available:
            result = run("watermark", [source], {"text": "SMOKE MARK"})[0]
            require(all("SMOKE MARK" in (page.extract_text() or "") for page in PdfReader(result).pages),
                    "watermark: mark missing on page")

        if "protect" in available:
            result = run("protect", [source], {"password": "smoke-password"})[0]
            reader = PdfReader(result)
            require(reader.is_encrypted, "protect: output is not encrypted")
            require(not reader.decrypt("wrong-password"), "protect: wrong password accepted")
            require(bool(reader.decrypt("smoke-password")), "protect: correct password rejected")
            require("Frozen worker page one" in reader.pages[0].extract_text(), "protect: content lost")
            encrypted = result

        if "unlock" in available:
            result = run("unlock", [encrypted], {"password": "smoke-password"})[0]
            require(not PdfReader(result).is_encrypted and "Frozen worker page two" in pdf_text(result),
                    "unlock: decryption roundtrip lost content")

        if "forms" in available:
            result = run("forms", [form], {"values": '{"nombre":"Ana"}'})[0]
            require(PdfReader(result).get_form_text_fields().get("nombre") == "Ana", "forms: field was not filled")

        if "jpg_to_pdf" in available:
            result = run("jpg_to_pdf", [picture])[0]
            reader = PdfReader(result)
            require(len(reader.pages) == 1 and bool(reader.pages[0].images), "jpg_to_pdf: image page missing")

        if "pdf_to_markdown" in available:
            result = run("pdf_to_markdown", [source])[0]
            require("Frozen worker page two" in result.read_text(encoding="utf-8"), "pdf_to_markdown: content missing")

        if "compress" in available:
            result = run("compress", [source])[0]
            require(result.stat().st_size <= source.stat().st_size and pdf_text(result) == pdf_text(source),
                    "compress: output enlarged or text changed")

        for path, before in originals.items():
            require(path.read_bytes() == before, f"Original modified: {path}")
    print(f"Packaged worker smoke passed: {', '.join(completed)}", flush=True)


if __name__ == "__main__":
    main()
