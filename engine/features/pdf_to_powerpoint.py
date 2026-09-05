from contextlib import closing
from pathlib import Path

def run(inputs: list[str], output_dir: str, options: dict) -> list[str]:
    if len(inputs) != 1:
        raise ValueError("Selecciona exactamente un archivo.")
    source = Path(inputs[0])
    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)

    import io
    import pypdfium2 as pdfium
    from pptx import Presentation
    from pptx.util import Inches
    presentation = Presentation()
    with pdfium.PdfDocument(str(source)) as pdf:
        if not len(pdf):
            raise ValueError("El PDF no tiene páginas.")
        with closing(pdf[0]) as first:
            presentation.slide_width = Inches(first.get_width() / 72)
            presentation.slide_height = Inches(first.get_height() / 72)
        for i in range(len(pdf)):
            with closing(pdf[i]) as page:
                if page.get_width() * page.get_height() * 4 > 100_000_000:
                    raise ValueError("Página demasiado grande para renderizar.")
                with closing(page.render(scale=2)) as bitmap:
                    image = bitmap.to_pil()
                    buffer = io.BytesIO()
                    image.save(buffer, format="PNG")
                    image.close()
                    buffer.seek(0)
                    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                    ratio = min(presentation.slide_width / page.get_width(), presentation.slide_height / page.get_height())
                    width, height = int(page.get_width() * ratio), int(page.get_height() * ratio)
                    slide.shapes.add_picture(buffer, (presentation.slide_width-width)//2, (presentation.slide_height-height)//2, width, height)
    target = out / "converted.pptx"
    presentation.save(target)
    return [str(target)]
