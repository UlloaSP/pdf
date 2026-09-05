import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
from engine.features import pdf_to_powerpoint as feature

class ConversionTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.root = Path(self.tmp.name)
        self.source = self.root / "sample.pdf"
        self.out = self.root / "output"
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(self.source))
        c.drawString(72, 720, "Hello PDF")
        c.showPage()
        c.drawString(72, 720, "Second page")
        c.save()

    def test_rejects_empty_input(self):
        with self.assertRaises(ValueError):
            feature.run([], str(self.out), {})

    def test_slide_per_page(self):
        from pptx import Presentation
        paths = feature.run([str(self.source)], str(self.out), {})
        slides = Presentation(paths[0]).slides
        self.assertEqual(len(slides), 2)
        self.assertEqual(len(slides[0].shapes), 1)
